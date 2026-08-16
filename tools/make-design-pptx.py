#!/usr/bin/env python3
"""Generate the robot design PowerPoint from the same content as
robot-design/chassis-designs-slides.html.

The HTML deck is the source of truth for wording. This script exists so the
.pptx can be regenerated rather than hand-edited into drift, and so the diagrams
land as real, editable PowerPoint shapes instead of flat screenshots.

Run:
    python3 tools/make-design-pptx.py

Writes robot-design/chassis-designs.pptx.

Font is Trebuchet MS on purpose: friendlier than Arial, and it ships on both
macOS and Windows, so nothing substitutes if the file gets emailed.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# --- BIOGLOW skin, matching the website and the HTML deck --------------------
INK    = RGBColor(0x16, 0x36, 0x2B)   # deep rainforest green
MUTED  = RGBColor(0x4A, 0x63, 0x57)
SOFT   = RGBColor(0x7D, 0x94, 0x88)
ACCENT = RGBColor(0xEB, 0x6C, 0x36)   # one bright signal; green-on-green has no contrast
PAPER  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xFC, 0xE6, 0xDC)   # accent at low opacity, flattened
GREY   = RGBColor(0xD8, 0xDE, 0xDA)

FONT = "Trebuchet MS"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = 13.333
H = 7.5


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    return s


def text(s, txt, x, y, w, h, size=18, color=MUTED, bold=False,
         align=PP_ALIGN.LEFT, italic=False, spacing=1.25):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = txt.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return box


def eyebrow(s, txt):
    text(s, txt.upper(), 0.9, 0.55, 8, 0.4, size=11, color=MUTED, bold=True)


def title(s, txt, size=40, y=1.05, color=INK):
    text(s, txt, 0.9, y, 11.5, 1.6, size=size, color=color, bold=True, spacing=1.05)


def lede(s, txt, y=2.35, size=15, w=10.6):
    text(s, txt, 0.9, y, w, 1.5, size=size, color=MUTED, spacing=1.35)


def aside(s, txt, y=6.35):
    """Editorial callout: accent rule on the left, italic text."""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(y),
                             Inches(0.035), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, txt, 1.1, y + 0.02, 10.6, 0.6, size=14, color=INK, italic=True)


def box(s, x, y, w, h, fill=WHITE, line=MUTED, lw=1.0, radius=0.04, rot=0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if rot:
        shp.rotation = rot
    shp.text_frame.text = ""
    return shp


def plain(s, x, y, w, h, fill=GREY):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def circle(s, cx, cy, r, fill=WHITE, line=MUTED, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                             Inches(2 * r), Inches(2 * r))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def line(s, x1, y1, x2, y2, color=MUTED, lw=1.25, dash=False, arrow=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        # python-pptx has no arrowhead API; drop into the XML.
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return c


def tag(s, txt, x, y, w, color=MUTED):
    """Small uppercase mono-ish label in a hairline box."""
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = txt.upper()
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = color
    return shp


def caption(s, txt, y=6.95):
    text(s, txt.upper(), 0.9, y, 11.5, 0.35, size=10, color=SOFT, bold=True)


# ============================================================== 1. TITLE
s = slide()
eyebrow(s, "BIOGLOW  ·  Robot design")
title(s, "Six teams built six robots.\nThe good ones all made the same choices.", size=38, y=1.5)
lede(s, "Before we build ours, let's look at what other people already worked out — "
        "and at the mistakes they only found after building.\n"
        "None of this is guesswork. It comes from a coach with ten years of seasons "
        "behind him, through regionals and nationals to an international final.", y=3.4)
aside(s, "The robot that wins is not the cleverest one. It is the one that does the "
         "same thing every single time.", y=5.4)

# ================================================ 2. ONE BASE, MANY ATTACHMENTS
s = slide()
eyebrow(s, "The big idea")
title(s, "One base. Many attachments.")
lede(s, "Do not build a robot per mission. Build one really good base that never changes, "
        "then build a small attachment for each job that clips onto it. The base is the "
        "thing you get reliable once; the attachments are where the cleverness goes.", y=2.15)

# base (focal) + three attachments, arrows fanned along the base edge
b = box(s, 1.6, 3.6, 2.9, 2.6, fill=TINT, line=ACCENT, lw=1.75)
tag(s, "Fixed", 1.85, 3.85, 0.95, color=ACCENT)
text(s, "One base", 1.6, 4.7, 2.9, 0.4, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Never changes", 1.6, 5.15, 2.9, 0.4, size=15, color=MUTED, align=PP_ALIGN.CENTER)

for i, label in enumerate(["Attachment for one job", "Attachment for another", "And another"]):
    ay = 3.6 + i * 0.92
    line(s, 4.5, ay + 0.36, 7.4, ay + 0.36, color=MUTED, arrow=True)
    ab = box(s, 7.4, ay, 3.9, 0.72, fill=WHITE, line=MUTED)
    text(s, label, 7.5, ay + 0.19, 3.7, 0.4, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)

caption(s, "Time is the thing you run out of. A base that never changes is a base you only debug once.")

# ============================================================ 3. THE LINE-UP
s = slide()
eyebrow(s, "The line-up")
title(s, "What each one is good at, and what it costs", size=34)
lede(s, "Read the right-hand column carefully. That is where the lesson is.", y=1.95, size=14)

rows = [
    ("Robot", "What it does well", "What it costs you"),
    ("HummerOne Pro",
     "Easiest to build. One kit plus the expansion set. Simple clip-on attachments.",
     "Wheels slip once an attachment gets heavy."),
    ("Xbot",
     "Frame right around the outside, so it balances well. Attachments on top and front.",
     "The frame eats a lot of LEGO. Wheels can still slip — blue SPIKE wheels."),
    ("YellowBot v1",
     "Box frame lines up against walls. Attachments drop in using gravity. Colour sensor at the front.",
     "Wheels can still slip — blue SPIKE wheels."),
    ("YellowBot v2",
     "Better box frame. Front-wheel drive. Hub sunk down low. Side walls can push things.",
     "Raised frame needs extra LEGO. Needs 62.4 mm wheels."),
    ("Track-X",
     "Sits very low, so it grips. Very accurate and repeatable. Fast frame swaps. Widely used.",
     "Needs 62.4 mm wheels. Open sides, so no walls to line up with."),
    ("Boxor",
     "Newest. Low, accurate, fast to swap, and a box frame for lining up. Best attachment system.",
     "Needs 62.4 mm wheels."),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.5),
                               Inches(11.5), Inches(3.9))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(5.1)
tbl.columns[2].width = Inches(4.4)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.42 if ri == 0 else 0.58)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.margin_left = cell.margin_right = Inches(0.09)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = INK
        elif rows[ri][0] == "Boxor":
            cell.fill.fore_color.rgb = TINT
        else:
            cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEF, 0xF2, 0xF0)
        p = cell.text_frame.paragraphs[0]
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(11 if ri else 11)
            r.font.bold = (ri == 0) or (ci == 0)
            if ri == 0:
                r.font.color.rgb = PAPER
            elif rows[ri][0] == "Boxor" and ci == 0:
                r.font.color.rgb = ACCENT
            else:
                r.font.color.rgb = INK if ci == 0 else MUTED

aside(s, "Boxor is not better because it is newer. It is better because it collected the "
         "good ideas from the five before it.", y=6.6)

# ======================================================== 4. THE WHEEL PATTERN
s = slide()
eyebrow(s, "Spot the pattern")
title(s, "Three robots slip. Three don't. Why?")
lede(s, "Sort the six by which wheels they use. The slipping stops exactly where the wheels "
        "change — and a coach who made that swap on a real team says the accuracy improved "
        "dramatically, with the robot no longer slipping on turns.", y=2.15)

# left: blue SPIKE
box(s, 0.9, 3.5, 5.0, 2.6, fill=RGBColor(0xEF, 0xF2, 0xF0), line=SOFT)
tag(s, "Blue SPIKE wheels", 1.15, 3.72, 2.3)
for i, n in enumerate(["HummerOne Pro", "Xbot", "YellowBot v1"]):
    text(s, n, 1.15, 4.25 + i * 0.38, 3.0, 0.35, size=15, color=INK, bold=True)
line(s, 1.15, 5.62, 3.1, 5.62, color=SOFT, lw=1.5)
circle(s, 2.75, 5.42, 0.2, line=MUTED)
line(s, 1.35, 5.62, 1.75, 5.62, color=MUTED, lw=2.5, dash=True)
text(s, "All three list slipping", 3.3, 5.3, 2.5, 0.6, size=15, color=MUTED, bold=True)

# right: 62.4 mm — focal
box(s, 6.45, 3.5, 5.0, 2.6, fill=TINT, line=ACCENT, lw=1.75)
tag(s, "62.4 mm wheels", 6.7, 3.72, 2.1, color=ACCENT)
for i, n in enumerate(["YellowBot v2", "Track-X", "Boxor"]):
    text(s, n, 6.7, 4.25 + i * 0.38, 3.0, 0.35, size=15, color=INK, bold=True)
line(s, 6.7, 5.62, 8.65, 5.62, color=SOFT, lw=1.5)
circle(s, 8.3, 5.42, 0.2, line=INK, lw=1.6)
line(s, 8.1, 5.62, 8.5, 5.62, color=INK, lw=3)
text(s, "None of them do", 8.85, 5.3, 2.5, 0.6, size=15, color=INK, bold=True)

# the swap arrow between the panels
line(s, 5.95, 4.8, 6.42, 4.8, color=ACCENT, lw=1.5, arrow=True)

caption(s, "The 62.4 mm wheel is also bigger, so it drives faster. More grip and more speed, same part.")

# ============================================================ 5. BOX FRAME
s = slide()
eyebrow(s, "Design rule 1")
title(s, "A box frame lines you up for free")
lede(s, "Push a box flat against the wall and it can only sit one way — square. That is a "
        "free, perfect starting position, every single run, with no code at all.", y=2.15)

plain(s, 0.9, 3.2, 10.5, 0.16, fill=RGBColor(0xC3, 0xCD, 0xC7))
text(s, "BORDER WALL", 0.9, 2.9, 3, 0.3, size=10, color=MUTED, bold=True)

# open sides, sitting crooked
box(s, 1.5, 3.42, 0.28, 1.5, fill=WHITE, line=MUTED, rot=-6)
box(s, 3.7, 3.42, 0.28, 1.5, fill=WHITE, line=MUTED, rot=-6)
box(s, 1.5, 4.66, 2.48, 0.28, fill=WHITE, line=MUTED, rot=-6)
text(s, "Open sides", 1.4, 5.3, 4.4, 0.35, size=17, color=INK, bold=True)
text(s, "Only one corner touches. It can sit\ncrooked and you will not notice.",
     1.4, 5.72, 4.6, 0.8, size=15, color=MUTED)

# box frame, flush — focal
box(s, 6.6, 3.36, 4.3, 1.6, fill=TINT, line=ACCENT, lw=1.75)
box(s, 7.05, 3.75, 3.4, 0.8, fill=WHITE, line=SOFT)
text(s, "HUB", 7.05, 4.0, 3.4, 0.3, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
plain(s, 6.6, 3.36, 4.3, 0.06, fill=ACCENT)
text(s, "Box frame", 6.6, 5.3, 4.6, 0.35, size=17, color=INK, bold=True)
text(s, "The whole edge touches, so it is square\nbefore it moves. Same start, every run.",
     6.6, 5.72, 4.8, 0.8, size=15, color=MUTED)

caption(s, "Track-X is fast and accurate but has open sides. Boxor kept the speed and added the frame.")

# ==================================================== 6. CENTRE OF GRAVITY
s = slide()
eyebrow(s, "Design rule 2")
title(s, "Keep the weight low — and even")
lede(s, "Weight up high makes the robot rock when it starts, stops and turns. Rocking lifts "
        "weight off the wheels, and a wheel with no weight on it spins instead of driving. "
        "Weight also needs spreading front to back: pile it all at one end and your turns "
        "stop being repeatable.", y=2.1)

GY = 5.75
line(s, 0.9, GY, 11.4, GY, color=SOFT, lw=1.75)
text(s, "MAT", 0.9, GY + 0.1, 2, 0.3, size=10, color=MUTED, bold=True)

# hub on top
box(s, 1.5, 5.0, 3.1, 0.6, fill=WHITE, line=MUTED)
box(s, 2.3, 4.05, 1.5, 0.95, fill=RGBColor(0xEF, 0xF2, 0xF0), line=MUTED)
text(s, "HUB", 2.3, 4.4, 1.5, 0.3, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
circle(s, 1.85, GY, 0.2, line=MUTED)
circle(s, 4.25, GY, 0.2, line=MUTED)
circle(s, 3.05, 4.9, 0.09, fill=MUTED, line=MUTED)
line(s, 5.0, 4.9, 5.0, GY, color=MUTED, lw=1, dash=True)
text(s, "Hub on top", 1.5, 3.45, 3.6, 0.35, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Weight sits high. It rocks.", 1.4, 3.82, 3.8, 0.3, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# hub recessed — focal
box(s, 7.0, 5.2, 3.7, 0.4, fill=TINT, line=ACCENT, lw=1.75)
box(s, 7.9, 5.28, 1.9, 0.26, fill=WHITE, line=SOFT)
circle(s, 7.35, GY, 0.2, line=INK, lw=1.6)
circle(s, 10.35, GY, 0.2, line=INK, lw=1.6)
circle(s, 8.85, 5.38, 0.09, fill=ACCENT, line=ACCENT)
line(s, 11.0, 5.38, 11.0, GY, color=MUTED, lw=1, dash=True)
text(s, "Hub sunk into the frame", 7.0, 3.45, 3.9, 0.35, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, "Weight stays on the wheels.", 7.0, 3.82, 3.9, 0.3, size=14, color=MUTED, align=PP_ALIGN.CENTER)

caption(s, "Dot = where the weight averages out. Lower is better, and centred is better still.")

# ======================================================= 7. WHEELS STRAIGHT
s = slide()
eyebrow(s, "Design rule 3")
title(s, "Hold the wheels straight")
lede(s, "Support the wheel on the motor side and on the far side. A heavy attachment — say "
        "you are carrying something across the table — presses down and bends the wheels "
        "slightly inward. You will not see it happen. You will only see the robot stop "
        "driving straight.", y=2.1)

# braced one side
plain(s, 1.1, 4.0, 0.3, 1.7, fill=RGBColor(0xEF, 0xF2, 0xF0))
line(s, 1.4, 4.85, 2.9, 4.85, color=MUTED, lw=2)
box(s, 2.9, 4.35, 0.38, 1.0, fill=WHITE, line=MUTED, lw=1.6, rot=7)
text(s, "Braced one side", 1.1, 3.5, 4.0, 0.35, size=17, color=INK, bold=True)
line(s, 3.8, 4.85, 5.4, 4.85, color=MUTED, lw=1.25, dash=True, arrow=True)
text(s, "DRIFTS", 4.1, 4.45, 1.4, 0.3, size=11, color=MUTED, bold=True)
text(s, "The axle flexes, the wheel leans, and every\nrun ends somewhere slightly different.",
     1.1, 5.95, 4.8, 0.8, size=15, color=MUTED)

# braced both sides — focal
plain(s, 6.7, 4.0, 0.3, 1.7, fill=RGBColor(0xEF, 0xF2, 0xF0))
plain(s, 9.0, 4.0, 0.3, 1.7, fill=RGBColor(0xEF, 0xF2, 0xF0))
line(s, 7.0, 4.85, 9.0, 4.85, color=INK, lw=2)
box(s, 7.8, 4.35, 0.38, 1.0, fill=TINT, line=ACCENT, lw=1.75)
text(s, "Braced both sides", 6.7, 3.5, 4.2, 0.35, size=17, color=INK, bold=True)
line(s, 9.7, 4.85, 11.3, 4.85, color=ACCENT, lw=1.25, arrow=True)
text(s, "STRAIGHT", 10.0, 4.45, 1.6, 0.3, size=11, color=ACCENT, bold=True)
text(s, "The wheel stays square, so the same\nprogram gives the same result.",
     6.7, 5.95, 4.8, 0.8, size=15, color=MUTED)

caption(s, "Top view · grey = frame beam")

# ========================================================== 8. ATTACHMENTS
s = slide()
eyebrow(s, "Design rule 4")
title(s, "Attachments should drop on, not screw on")
lede(s, "In a match you get two and a half minutes. Every second spent wrestling an "
        "attachment is a second not scoring. This is the gravity technique: a couple of "
        "small pegs guide it into place, then its own weight holds it. Nothing locks.", y=2.1)

box(s, 4.7, 3.35, 3.5, 0.8, fill=TINT, line=ACCENT, lw=1.75)
text(s, "Attachment", 4.7, 3.58, 3.5, 0.35, size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
line(s, 6.45, 4.25, 6.45, 4.95, color=ACCENT, lw=1.5, arrow=True)
text(s, "DROPS IN", 6.65, 4.45, 1.6, 0.3, size=11, color=ACCENT, bold=True)

box(s, 4.3, 5.05, 4.3, 1.5, fill=WHITE, line=INK, lw=1.6)
text(s, "Chassis", 4.3, 5.65, 4.3, 0.35, size=16, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
plain(s, 5.3, 4.97, 2.3, 0.16)
text(s, "TOP MOUNT", 8.75, 4.92, 2.2, 0.3, size=10, color=MUTED, bold=True)
plain(s, 8.52, 5.4, 0.16, 0.9)
text(s, "FRONT MOUNT", 8.75, 5.72, 2.4, 0.3, size=10, color=MUTED, bold=True)
circle(s, 4.6, 6.6, 0.17, line=MUTED)
circle(s, 8.3, 6.6, 0.17, line=MUTED)

text(s, "Held by its own weight", 0.9, 3.5, 3.2, 0.6, size=16, color=INK, bold=True)
text(s, "No pins to line up\nunder pressure.", 0.9, 4.15, 3.2, 0.7, size=15, color=MUTED)

# ======================================================= 9. SENSORS AND CODE
s = slide()
eyebrow(s, "Design rule 5")
title(s, "The gyro does the driving, not the motor blocks")
lede(s, "The ready-made movement blocks look easy and they will let you down. Zero the gyro "
        "at the start, then hold that heading the whole way, correcting continuously. Bump a "
        "model mid-run and it steers itself back to zero.", y=2.1)

items = [
    ("Zero the gyro, then hold the heading.",
     "Drive at 0°, measure the error constantly, steer by an amount proportional to how far off you are."),
    ("Motor encoders for distance.",
     "Exact degrees forward, exact degrees to turn. Repeatable in a way that timing never is."),
    ("One colour sensor is still worth having.",
     "But black lines are disappearing from FIRST's mats, so most teams now navigate on gyro and encoders."),
    ("Four wheels with rubber tyres.",
     "Not a ball-bearing castor. Four driven wheels track straighter."),
]
y = 3.5
for head, sub in items:
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y + 0.06),
                            Inches(0.22), Inches(0.22))
    bx.adjustments[0] = 0.2
    bx.fill.background(); bx.line.color.rgb = ACCENT; bx.line.width = Pt(1.5)
    bx.shadow.inherit = False
    text(s, head, 1.3, y, 10.0, 0.32, size=16, color=INK, bold=True)
    text(s, sub, 1.3, y + 0.36, 10.0, 0.32, size=13, color=MUTED)
    y += 0.85

aside(s, "We already have this: the Acceleration block ramps up, holds the gyro heading and "
         "eases off — and the same maths is drive_cm_gyro in the Python toolkit.", y=6.72)

# ============================================================ 10. OUR TURN
s = slide()
eyebrow(s, "Our turn")
title(s, "What we are deciding this week")
lede(s, "We do not have to invent a robot. We have to make these choices on purpose, and be "
        "able to say why — because that is the question the judges ask.", y=2.05)

checks = [
    "Is the base finished, and can we leave it alone?",
    "Box frame or open sides?",
    "Which wheels?",
    "Where does the hub sit, and is the weight even?",
    "Are the wheels supported on both sides?",
    "How fast can we swap an attachment?",
    "Top mount, front mount, or both?",
    "Gyro and encoders, or the ready-made blocks?",
]
y = 3.0
for c in checks:
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y + 0.03),
                            Inches(0.24), Inches(0.24))
    bx.adjustments[0] = 0.2
    bx.fill.background(); bx.line.color.rgb = ACCENT; bx.line.width = Pt(1.5)
    bx.shadow.inherit = False
    text(s, c, 1.35, y, 10.5, 0.32, size=16, color=INK)
    y += 0.44

aside(s, "Write down what we choose and why, and photograph version one before we take it "
         "apart — that record is what the judges score hardest.", y=6.75)

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "robot-design", "chassis-designs.pptx")
prs.save(out)
print("wrote", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
